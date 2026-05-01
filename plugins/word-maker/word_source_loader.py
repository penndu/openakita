"""Source file extraction for word-maker projects."""

from __future__ import annotations

import csv
import json
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any

TEXT_SUFFIXES = {".md", ".txt", ".rst", ".log"}
CSV_SUFFIXES = {".csv"}
JSON_SUFFIXES = {".json"}


@dataclass(slots=True)
class SourceLoadResult:
    path: str
    source_type: str
    text: str
    ok: bool = True
    error: str = ""
    metadata: dict[str, Any] | None = None

    def to_dict(self) -> dict[str, Any]:
        return asdict(self)


def _truncate(text: str, max_chars: int) -> str:
    return text if len(text) <= max_chars else text[:max_chars] + "\n...[truncated]"


def _read_text(path: Path, max_chars: int) -> str:
    return _truncate(path.read_text(encoding="utf-8", errors="replace"), max_chars)


def _read_csv(path: Path, max_chars: int) -> str:
    lines: list[str] = []
    with path.open("r", encoding="utf-8", errors="replace", newline="") as handle:
        reader = csv.reader(handle)
        for index, row in enumerate(reader):
            lines.append(" | ".join(row))
            if index >= 50:
                lines.append("...[truncated rows]")
                break
    return _truncate("\n".join(lines), max_chars)


def _read_json(path: Path, max_chars: int) -> str:
    raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        parsed = json.loads(raw)
    except json.JSONDecodeError:
        return _truncate(raw, max_chars)
    return _truncate(json.dumps(parsed, ensure_ascii=False, indent=2), max_chars)


def _read_docx(path: Path, max_chars: int) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required to extract DOCX files") from exc

    document = Document(str(path))
    chunks: list[str] = []
    chunks.extend(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                chunks.append(" | ".join(cells))
    return _truncate("\n".join(chunks), max_chars)


def _read_xlsx(path: Path, max_chars: int) -> str:
    try:
        import openpyxl
    except ImportError as exc:
        raise RuntimeError("openpyxl is required to extract XLSX files") from exc

    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    chunks: list[str] = []
    try:
        for sheet in workbook.worksheets[:5]:
            chunks.append(f"# Sheet: {sheet.title}")
            for row_index, row in enumerate(sheet.iter_rows(values_only=True)):
                values = ["" if value is None else str(value) for value in row]
                if any(values):
                    chunks.append(" | ".join(values))
                if row_index >= 50:
                    chunks.append("...[truncated rows]")
                    break
    finally:
        workbook.close()
    return _truncate("\n".join(chunks), max_chars)


def _read_pptx(path: Path, max_chars: int) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required to extract PPTX files") from exc

    presentation = Presentation(str(path))
    chunks: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        chunks.append(f"# Slide {slide_index}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text.strip():
                chunks.append(text.strip())
    return _truncate("\n".join(chunks), max_chars)


def _read_pdf(path: Path, max_chars: int) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("pypdf is required to extract PDF files") from exc

    reader = PdfReader(str(path))
    chunks = [(page.extract_text() or "").strip() for page in reader.pages[:20]]
    return _truncate("\n\n".join(chunk for chunk in chunks if chunk), max_chars)


def load_source(path: str | Path, *, max_chars: int = 20000) -> SourceLoadResult:
    target = Path(path)
    suffix = target.suffix.lower()
    if not target.exists() or not target.is_file():
        return SourceLoadResult(str(target), "unknown", "", ok=False, error="File not found")

    try:
        if suffix in TEXT_SUFFIXES:
            text = _read_text(target, max_chars)
            source_type = suffix.lstrip(".") or "text"
        elif suffix in CSV_SUFFIXES:
            text = _read_csv(target, max_chars)
            source_type = "csv"
        elif suffix in JSON_SUFFIXES:
            text = _read_json(target, max_chars)
            source_type = "json"
        elif suffix == ".docx":
            text = _read_docx(target, max_chars)
            source_type = "docx"
        elif suffix == ".xlsx":
            text = _read_xlsx(target, max_chars)
            source_type = "xlsx"
        elif suffix == ".pptx":
            text = _read_pptx(target, max_chars)
            source_type = "pptx"
        elif suffix == ".pdf":
            text = _read_pdf(target, max_chars)
            source_type = "pdf"
        else:
            return SourceLoadResult(
                str(target),
                "unsupported",
                "",
                ok=False,
                error=f"Unsupported source file type: {suffix or '(none)'}",
            )
    except Exception as exc:
        return SourceLoadResult(str(target), suffix.lstrip(".") or "unknown", "", ok=False, error=str(exc))

    return SourceLoadResult(
        str(target),
        source_type,
        text,
        ok=True,
        metadata={"chars": len(text), "filename": target.name},
    )


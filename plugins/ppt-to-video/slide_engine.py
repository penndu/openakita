"""ppt-to-video — pure-logic core.

Pipeline (one PPTX in → one MP4 out):

1. ``soffice --headless --convert-to png`` exports each slide as a PNG.
2. ``python-pptx`` extracts the speaker-notes paragraph for each slide.
3. For every slide we synthesize the notes via the chosen TTS provider
   (reusing ``plugins/avatar-speaker/providers.py``); slides with empty
   notes get a configurable silent gap so the visual still gets a
   beat of screen time.
4. Each (image, audio) pair is encoded as a short clip with
   ``ffmpeg -loop 1 -i image -i audio -shortest`` (or ``-t silence_sec``
   when there is no audio).
5. The clips are concatenated with the ffmpeg concat demuxer into the
   final MP4.

Heavy / non-deterministic dependencies are deliberately *lazy*:

* ``python-pptx`` — only imported inside ``extract_slide_notes``.
* The TTS provider — only loaded by the worker; planning never touches
  network / models.
* ``soffice`` and ``ffmpeg`` — discovered via ``shutil.which`` so the
  plugin module imports cleanly on machines that don't have them, and
  ``check_deps`` can serve install guidance instead of crashing.

Keeping this layer thin (no FastAPI, no asyncio, no sqlite) means the
same engine can be reused inside ``shorts-batch`` (D3, future Sprint
17) when an LLM wants to assemble a slide-style explainer.
"""

from __future__ import annotations

import importlib
import shutil
import subprocess
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable, Iterable

from openakita_plugin_sdk.contrib import (
    KIND_NUMBER,
    KIND_OTHER,
    LowConfidenceField,
    Verification,
)

__all__ = [
    "DEFAULT_FPS",
    "DEFAULT_SILENT_SLIDE_SEC",
    "DEFAULT_VOICE",
    "DEFAULT_TTS_PROVIDER",
    "DEFAULT_RENDER_TIMEOUT_SEC",
    "SUPPORTED_INPUT_EXTENSIONS",
    "SlideMeta",
    "SlidePlan",
    "SlideVideoResult",
    "build_concat_command",
    "build_image_clip_command",
    "extract_slide_notes",
    "ffmpeg_available",
    "libreoffice_available",
    "plan_video",
    "pptx_available",
    "render_clips",
    "resolve_libreoffice",
    "run_pipeline",
    "soffice_convert_command",
    "to_verification",
]


# ── Constants ──────────────────────────────────────────────────────────


DEFAULT_FPS = 25
# Slides whose speaker notes are empty get this many seconds of screen
# time so the viewer at least sees the visual.  Audited keeps it short
# (2.0s) — the SKILL.md tip is "leave notes for every slide".
DEFAULT_SILENT_SLIDE_SEC = 2.0
DEFAULT_VOICE = "zh-CN-XiaoxiaoNeural"
DEFAULT_TTS_PROVIDER = "auto"
DEFAULT_RENDER_TIMEOUT_SEC = 1800.0
SUPPORTED_INPUT_EXTENSIONS = (".pptx", ".ppt", ".odp")


# ── Models ─────────────────────────────────────────────────────────────


@dataclass
class SlideMeta:
    """One slide after planning (image known, audio not yet rendered)."""

    index: int            # 1-based
    image_path: str
    notes: str
    audio_path: str | None = None
    audio_duration_sec: float = 0.0
    clip_path: str | None = None

    def to_dict(self) -> dict[str, Any]:
        return {
            "index": self.index,
            "image_path": self.image_path,
            "notes": self.notes,
            "audio_path": self.audio_path,
            "audio_duration_sec": self.audio_duration_sec,
            "clip_path": self.clip_path,
        }


@dataclass
class SlidePlan:
    """Frozen description of what the worker will produce."""

    input_path: str
    output_path: str
    work_dir: str
    voice: str
    tts_provider: str
    silent_slide_sec: float
    fps: int
    crf: int
    libx264_preset: str
    slides: list[SlideMeta] = field(default_factory=list)

    @property
    def slide_count(self) -> int:
        return len(self.slides)

    @property
    def notes_total_chars(self) -> int:
        return sum(len(s.notes) for s in self.slides)

    @property
    def empty_notes_count(self) -> int:
        return sum(1 for s in self.slides if not s.notes.strip())

    def to_dict(self) -> dict[str, Any]:
        return {
            "input_path": self.input_path,
            "output_path": self.output_path,
            "work_dir": self.work_dir,
            "voice": self.voice,
            "tts_provider": self.tts_provider,
            "silent_slide_sec": self.silent_slide_sec,
            "fps": self.fps,
            "crf": self.crf,
            "libx264_preset": self.libx264_preset,
            "slide_count": self.slide_count,
            "notes_total_chars": self.notes_total_chars,
            "empty_notes_count": self.empty_notes_count,
            "slides": [s.to_dict() for s in self.slides],
        }


@dataclass
class SlideVideoResult:
    """What the worker produced."""

    plan: SlidePlan
    output_path: str
    elapsed_sec: float
    slide_count: int
    audio_total_sec: float
    output_size_bytes: int
    tts_provider_used: str
    tts_fallbacks: int = 0  # how many slides fell back to silent gap

    def to_dict(self) -> dict[str, Any]:
        return {
            "plan": self.plan.to_dict(),
            "output_path": self.output_path,
            "elapsed_sec": self.elapsed_sec,
            "slide_count": self.slide_count,
            "audio_total_sec": self.audio_total_sec,
            "output_size_bytes": self.output_size_bytes,
            "tts_provider_used": self.tts_provider_used,
            "tts_fallbacks": self.tts_fallbacks,
        }


# ── Dep helpers ────────────────────────────────────────────────────────


def libreoffice_available() -> bool:
    """``True`` iff ``soffice`` is discoverable on PATH or in well-known dirs."""
    return resolve_libreoffice() is not None


def resolve_libreoffice() -> str | None:
    """Return the absolute path to a usable ``soffice`` binary, or ``None``.

    Order: ``$PATH`` (Linux/macOS/Windows-with-PATH), then the two
    standard Windows install dirs.  Returning a string (instead of just
    a bool) lets the worker pass the absolute path to ``subprocess`` so
    we don't have to depend on ``$PATH`` being correctly forwarded by
    the host process.
    """
    via_path = shutil.which("soffice")
    if via_path:
        return via_path
    for candidate in (
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/usr/bin/soffice",
        "/usr/local/bin/soffice",
    ):
        if Path(candidate).is_file():
            return candidate
    return None


def pptx_available() -> bool:
    """``True`` iff the ``python-pptx`` package can be imported."""
    return importlib.util.find_spec("pptx") is not None


def ffmpeg_available() -> bool:
    return shutil.which("ffmpeg") is not None


# ── PPTX → notes ───────────────────────────────────────────────────────


def extract_slide_notes(pptx_path: str | Path) -> list[str]:
    """Return one notes string per slide (empty string for no notes).

    Uses ``python-pptx`` for parsing.  We always return a list whose
    length equals the slide count even when notes are missing — callers
    rely on the index alignment with the exported PNGs.
    """
    if not pptx_available():
        raise ImportError(
            "python-pptx is not installed. Install with `pip install python-pptx`.",
        )
    if not Path(pptx_path).is_file():
        raise FileNotFoundError(f"pptx file not found: {pptx_path}")
    from pptx import Presentation  # local import — see module docstring

    prs = Presentation(str(pptx_path))
    out: list[str] = []
    for slide in prs.slides:
        notes_text = ""
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            if tf is not None:
                notes_text = (tf.text or "").strip()
        out.append(notes_text)
    return out


# ── PPTX → PNG via LibreOffice ─────────────────────────────────────────


def soffice_convert_command(
    *, soffice: str, input_path: str | Path, out_dir: str | Path,
    image_format: str = "png",
) -> list[str]:
    """Build the ``soffice`` argv that exports a presentation to images.

    LibreOffice's ``--convert-to png`` exports **one PNG per slide** when
    given a presentation file; the result is dropped into ``out_dir``.

    The command is returned (not executed) so callers can stub it in
    tests and so the exact argv ends up in logs / verification output.
    """
    Path(out_dir).mkdir(parents=True, exist_ok=True)
    return [
        str(soffice),
        "--headless",
        "--convert-to", image_format,
        "--outdir", str(out_dir),
        str(input_path),
    ]


def discover_exported_pngs(out_dir: str | Path, *, prefix: str | None = None) -> list[Path]:
    """List PNGs LibreOffice produced, sorted by slide index when possible.

    LibreOffice names the output ``<basename>.png`` for a single-slide
    deck, but for multi-slide decks newer versions emit
    ``<basename>-<N>.png``.  We sort by the trailing number when one is
    present, otherwise by filename — this keeps slide order stable.
    """
    d = Path(out_dir)
    if not d.is_dir():
        return []
    pngs = sorted(d.glob("*.png"))
    if prefix:
        pngs = [p for p in pngs if p.stem.startswith(prefix)]

    def _slide_idx(p: Path) -> tuple[int, str]:
        # last "-NNN" suffix wins; fall back to alphabetical
        parts = p.stem.rsplit("-", 1)
        if len(parts) == 2 and parts[1].isdigit():
            return (int(parts[1]), p.stem)
        return (10**9, p.stem)  # unindexed → after indexed

    return sorted(pngs, key=_slide_idx)


# ── Per-slide clip + concat commands ───────────────────────────────────


def build_image_clip_command(
    *, image_path: str | Path, audio_path: str | Path | None,
    duration_sec: float, output_path: str | Path,
    fps: int = DEFAULT_FPS, crf: int = 20, libx264_preset: str = "fast",
    ffmpeg: str = "ffmpeg",
) -> list[str]:
    """Build the ffmpeg argv that turns one image (+optional audio) into a clip.

    Three modes:

    * ``audio_path`` is given → ``-loop 1 -i image -i audio -shortest``
      (output length matches the audio).
    * ``audio_path`` is ``None`` and ``duration_sec`` > 0 → ``-loop 1 -i
      image -t duration_sec`` with a silent ``anullsrc`` audio so the
      concat-demuxer doesn't choke on missing audio streams.
    * Both unset → raises ``ValueError`` (caller bug).
    """
    if audio_path is None and duration_sec <= 0:
        raise ValueError(
            "build_image_clip_command requires audio_path or duration_sec > 0",
        )
    bin_path = ffmpeg if Path(ffmpeg).is_absolute() else (shutil.which(ffmpeg) or ffmpeg)
    common_video = [
        "-c:v", "libx264",
        "-preset", str(libx264_preset),
        "-crf", str(crf),
        "-pix_fmt", "yuv420p",
        "-r", str(fps),
        # libx264 needs even dimensions; pad if necessary so portrait
        # screenshots from soffice don't blow up the encoder.
        "-vf", "pad=ceil(iw/2)*2:ceil(ih/2)*2",
    ]

    if audio_path is not None:
        return [
            bin_path, "-y", "-hide_banner", "-loglevel", "error",
            "-loop", "1", "-i", str(image_path),
            "-i", str(audio_path),
            *common_video,
            "-c:a", "aac", "-b:a", "128k",
            "-shortest",
            str(output_path),
        ]

    return [
        bin_path, "-y", "-hide_banner", "-loglevel", "error",
        "-loop", "1", "-i", str(image_path),
        "-f", "lavfi", "-t", f"{duration_sec:.3f}",
        "-i", "anullsrc=channel_layout=stereo:sample_rate=44100",
        *common_video,
        "-c:a", "aac", "-b:a", "128k",
        "-t", f"{duration_sec:.3f}",
        str(output_path),
    ]


def build_concat_command(
    *, clip_paths: Iterable[str | Path], list_file: str | Path,
    output_path: str | Path, ffmpeg: str = "ffmpeg",
) -> list[str]:
    """Build the ffmpeg concat-demuxer argv + write the auxiliary list file.

    The concat demuxer reads a manifest like ``file 'clip1.mp4'`` and
    stitches the clips without re-encoding (``-c copy``), so total
    runtime stays linear in the number of slides instead of quadratic.
    """
    bin_path = ffmpeg if Path(ffmpeg).is_absolute() else (shutil.which(ffmpeg) or ffmpeg)
    paths = [Path(p) for p in clip_paths]
    if not paths:
        raise ValueError("build_concat_command requires at least one clip")
    list_p = Path(list_file)
    list_p.parent.mkdir(parents=True, exist_ok=True)
    list_p.write_text(
        "\n".join(f"file '{p.as_posix()}'" for p in paths) + "\n",
        encoding="utf-8",
    )
    return [
        bin_path, "-y", "-hide_banner", "-loglevel", "error",
        "-f", "concat", "-safe", "0",
        "-i", str(list_p),
        "-c", "copy",
        "-movflags", "+faststart",
        str(output_path),
    ]


# ── Plan & run ─────────────────────────────────────────────────────────


def plan_video(
    *,
    input_path: str,
    output_path: str,
    work_dir: str,
    voice: str = DEFAULT_VOICE,
    tts_provider: str = DEFAULT_TTS_PROVIDER,
    silent_slide_sec: float = DEFAULT_SILENT_SLIDE_SEC,
    fps: int = DEFAULT_FPS,
    crf: int = 20,
    libx264_preset: str = "fast",
    convert_runner: Callable[[list[str]], None] | None = None,
    notes_extractor: Callable[[str], list[str]] | None = None,
    soffice_path: str | None = None,
    pngs_discoverer: Callable[[str], list[Path]] | None = None,
) -> SlidePlan:
    """Validate inputs, run soffice, extract notes, return a frozen plan.

    Side effects: writes PNG files into ``work_dir`` (or whatever
    ``convert_runner`` decides).  Network is never touched here — TTS
    happens in :func:`run_pipeline`.

    The four ``*_runner`` / ``*_extractor`` parameters are dependency
    injection points so tests can avoid spawning soffice / installing
    python-pptx.  The defaults call the real implementations.
    """
    if not input_path or not str(input_path).strip():
        raise ValueError("input_path is required and must not be empty")
    if not output_path or not str(output_path).strip():
        raise ValueError("output_path is required and must not be empty")
    if not str(output_path).lower().endswith(".mp4"):
        raise ValueError(
            f"output_path must end in .mp4 (got {output_path!r})",
        )
    if not (0.5 <= silent_slide_sec <= 30.0):
        raise ValueError(
            f"silent_slide_sec must be in [0.5, 30.0], got {silent_slide_sec!r}",
        )
    ext = Path(input_path).suffix.lower()
    if ext not in SUPPORTED_INPUT_EXTENSIONS:
        raise ValueError(
            f"unsupported input extension {ext!r}; "
            f"supported: {SUPPORTED_INPUT_EXTENSIONS}",
        )
    if not Path(input_path).is_file():
        raise FileNotFoundError(f"input file not found: {input_path}")

    work_dir_p = Path(work_dir)
    images_dir = work_dir_p / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    # 1. soffice convert — runner can be stubbed in tests
    if convert_runner is None:
        soffice = soffice_path or resolve_libreoffice()
        if not soffice:
            raise FileNotFoundError(
                "LibreOffice (soffice) was not found. Install LibreOffice "
                "(https://www.libreoffice.org) or pass soffice_path=...",
            )
        cmd = soffice_convert_command(
            soffice=soffice, input_path=input_path, out_dir=images_dir,
        )

        def _default_runner(c: list[str]) -> None:
            subprocess.run(c, check=True, timeout=DEFAULT_RENDER_TIMEOUT_SEC)

        _default_runner(cmd)
    else:
        cmd = soffice_convert_command(
            soffice=soffice_path or "soffice",
            input_path=input_path, out_dir=images_dir,
        )
        convert_runner(cmd)

    # 2. discover PNGs
    if pngs_discoverer is not None:
        pngs = list(pngs_discoverer(str(images_dir)))
    else:
        pngs = discover_exported_pngs(images_dir)
    if not pngs:
        raise RuntimeError(
            f"LibreOffice did not produce any PNGs in {images_dir}; "
            "the source file may be empty or corrupted.",
        )

    # 3. extract notes (length-aligned with slide count)
    extractor = notes_extractor or extract_slide_notes
    notes = list(extractor(input_path))
    # If notes count mismatches png count, pad / truncate to PNG count.
    # Notes drive narration, but PNG order drives the visual timeline,
    # so PNGs win when there's a discrepancy.
    if len(notes) < len(pngs):
        notes = notes + [""] * (len(pngs) - len(notes))
    elif len(notes) > len(pngs):
        notes = notes[: len(pngs)]

    slides = [
        SlideMeta(index=i + 1, image_path=str(pngs[i]), notes=notes[i])
        for i in range(len(pngs))
    ]

    return SlidePlan(
        input_path=str(input_path),
        output_path=str(output_path),
        work_dir=str(work_dir_p),
        voice=str(voice),
        tts_provider=str(tts_provider),
        silent_slide_sec=float(silent_slide_sec),
        fps=int(fps),
        crf=int(crf),
        libx264_preset=str(libx264_preset),
        slides=slides,
    )


def render_clips(
    plan: SlidePlan, *, ffmpeg_runner: Callable[[list[str]], None],
    on_progress: Callable[[int, int], None] | None = None,
) -> tuple[float, int]:
    """Render one ffmpeg clip per slide and return (audio_total, fallback_count).

    Mutates each :class:`SlideMeta` in ``plan.slides`` to set
    ``clip_path`` and (when no audio was provided) records a default
    ``audio_duration_sec`` of ``plan.silent_slide_sec``.  ``ffmpeg_runner``
    is injected so tests can stub it.
    """
    audio_total = 0.0
    fallback_count = 0
    clips_dir = Path(plan.work_dir) / "clips"
    clips_dir.mkdir(parents=True, exist_ok=True)

    for i, slide in enumerate(plan.slides):
        clip_p = clips_dir / f"slide_{slide.index:04d}.mp4"
        if slide.audio_path:
            cmd = build_image_clip_command(
                image_path=slide.image_path,
                audio_path=slide.audio_path,
                duration_sec=slide.audio_duration_sec,
                output_path=clip_p,
                fps=plan.fps, crf=plan.crf,
                libx264_preset=plan.libx264_preset,
            )
            audio_total += slide.audio_duration_sec
        else:
            duration = plan.silent_slide_sec
            slide.audio_duration_sec = duration
            cmd = build_image_clip_command(
                image_path=slide.image_path,
                audio_path=None, duration_sec=duration,
                output_path=clip_p,
                fps=plan.fps, crf=plan.crf,
                libx264_preset=plan.libx264_preset,
            )
            fallback_count += 1
            audio_total += duration
        ffmpeg_runner(cmd)
        slide.clip_path = str(clip_p)
        if on_progress is not None:
            on_progress(i + 1, len(plan.slides))

    return audio_total, fallback_count


def run_pipeline(
    plan: SlidePlan,
    *,
    tts_synth: Callable[[str, str], tuple[Path, float]] | None = None,
    ffmpeg_runner: Callable[[list[str]], None] | None = None,
    on_progress: Callable[[str, int, int], None] | None = None,
) -> SlideVideoResult:
    """Synthesize TTS, render per-slide clips, concat into the final MP4.

    Args:
        plan: Output of :func:`plan_video`.
        tts_synth: ``Callable[(text, voice), (audio_path, duration_sec)]``.
            Tests pass a fake; the worker passes a wrapper around
            avatar-speaker's ``synthesize`` coroutine.  Returning
            ``None`` from this callable signals "no audio for this slide
            → use silent gap".
        ffmpeg_runner: ``Callable[[list[str]], None]``.  Defaults to
            ``subprocess.run(check=True, timeout=...)``.
        on_progress: ``Callable[[stage, done, total], None]``.

    The function is synchronous; the plugin's worker runs it inside
    ``asyncio.to_thread`` so the event loop stays responsive.
    """
    started = time.monotonic()
    ffmpeg_runner = ffmpeg_runner or _default_ffmpeg_runner
    audio_dir = Path(plan.work_dir) / "audio"
    audio_dir.mkdir(parents=True, exist_ok=True)
    tts_provider_used = "none"
    tts_fallbacks = 0

    # 1. TTS per slide.
    if tts_synth is not None:
        for i, slide in enumerate(plan.slides):
            text = slide.notes.strip()
            if not text:
                tts_fallbacks += 1
                if on_progress:
                    on_progress("tts", i + 1, len(plan.slides))
                continue
            try:
                result = tts_synth(text, plan.voice)
            except Exception:  # noqa: BLE001 — fall back per slide, never fail the whole job
                tts_fallbacks += 1
                if on_progress:
                    on_progress("tts", i + 1, len(plan.slides))
                continue
            if result is None:
                tts_fallbacks += 1
            else:
                audio_path, duration_sec = result
                slide.audio_path = str(audio_path)
                slide.audio_duration_sec = float(duration_sec)
                tts_provider_used = "injected"
            if on_progress:
                on_progress("tts", i + 1, len(plan.slides))
    else:
        tts_fallbacks = len(plan.slides)

    # 2. Per-slide clips.
    audio_total, render_fallbacks = render_clips(
        plan,
        ffmpeg_runner=ffmpeg_runner,
        on_progress=lambda d, t: on_progress("clips", d, t) if on_progress else None,
    )
    # render_clips' own fallback counter only counts slides where
    # tts_synth returned None *and* the worker fell back to silence
    # — same number we already track in ``tts_fallbacks``.
    tts_fallbacks = max(tts_fallbacks, render_fallbacks)

    # 3. Concat.
    list_file = Path(plan.work_dir) / "concat.txt"
    cmd = build_concat_command(
        clip_paths=[s.clip_path for s in plan.slides if s.clip_path],
        list_file=list_file,
        output_path=plan.output_path,
    )
    ffmpeg_runner(cmd)
    if on_progress:
        on_progress("concat", 1, 1)

    elapsed = time.monotonic() - started
    out_p = Path(plan.output_path)
    size_bytes = out_p.stat().st_size if out_p.is_file() else 0

    return SlideVideoResult(
        plan=plan,
        output_path=str(out_p),
        elapsed_sec=elapsed,
        slide_count=len(plan.slides),
        audio_total_sec=audio_total,
        output_size_bytes=size_bytes,
        tts_provider_used=tts_provider_used,
        tts_fallbacks=tts_fallbacks,
    )


def _default_ffmpeg_runner(cmd: list[str]) -> None:
    subprocess.run(cmd, check=True, timeout=DEFAULT_RENDER_TIMEOUT_SEC)


# ── Verification (D2.10) ───────────────────────────────────────────────


def to_verification(result: SlideVideoResult) -> Verification:
    """Convert a :class:`SlideVideoResult` into a D2.10 verification envelope.

    Yellow flags (the output still ships, but humans should glance):
      * zero slides produced (likely soffice failed silently),
      * majority of slides have empty notes (the video will feel mute),
      * TTS fell back to silence for >50% of slides,
      * output file is 0 bytes.
    """
    fields: list[LowConfidenceField] = []

    if result.slide_count == 0:
        fields.append(LowConfidenceField(
            path="$.slide_count",
            value=0,
            kind=KIND_NUMBER,
            reason="no slides were rendered — LibreOffice probably "
                   "failed to convert the input; check check-deps and "
                   "the source file integrity",
        ))

    empty_ratio = (
        result.plan.empty_notes_count / result.slide_count
        if result.slide_count else 0.0
    )
    if result.slide_count > 0 and empty_ratio > 0.5:
        fields.append(LowConfidenceField(
            path="$.plan.empty_notes_count",
            value=result.plan.empty_notes_count,
            kind=KIND_NUMBER,
            reason=(
                f"{result.plan.empty_notes_count}/{result.slide_count} slides "
                "have empty speaker notes — narration will rely on the "
                "configured silent gap, which may feel like a mute video"
            ),
        ))

    if result.slide_count > 0:
        fb_ratio = result.tts_fallbacks / result.slide_count
        if fb_ratio > 0.5 and result.tts_fallbacks > result.plan.empty_notes_count:
            fields.append(LowConfidenceField(
                path="$.tts_fallbacks",
                value=result.tts_fallbacks,
                kind=KIND_NUMBER,
                reason=(
                    f"TTS fell back to silence for {result.tts_fallbacks}/"
                    f"{result.slide_count} slides — the chosen provider "
                    "may be misconfigured or rate-limited"
                ),
            ))

    if result.output_size_bytes == 0:
        fields.append(LowConfidenceField(
            path="$.output_size_bytes",
            value=0,
            kind=KIND_NUMBER,
            reason="output file is 0 bytes — the ffmpeg concat probably "
                   "failed silently; check disk space and codecs",
        ))

    if result.plan.tts_provider == "stub":
        fields.append(LowConfidenceField(
            path="$.plan.tts_provider",
            value=result.plan.tts_provider,
            kind=KIND_OTHER,
            reason="TTS provider is the stub (silent placeholder) — "
                   "configure edge / dashscope / openai for real narration",
        ))

    return Verification(
        verified=not fields,
        verifier_id="ppt_to_video_self_check",
        low_confidence_fields=fields,
    )

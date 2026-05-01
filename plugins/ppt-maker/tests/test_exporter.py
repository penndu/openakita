from __future__ import annotations

import pytest
from ppt_design import DesignBuilder
from ppt_exporter import PptxExporter
from ppt_ir import SlideIrBuilder
from ppt_layouts import LAYOUT_REGISTRY
from ppt_models import DeckMode, SlideType
from ppt_outline import OutlineBuilder
from pptx import Presentation


def test_exporter_writes_editable_pptx(tmp_path) -> None:
    outline = OutlineBuilder().build(mode=DeckMode.TOPIC_TO_DECK, title="Roadmap", slide_count=3)
    design = DesignBuilder().build(outline=outline)
    ir = SlideIrBuilder().build(outline=outline, spec_lock=design["spec_lock"])

    path = PptxExporter().export(ir, tmp_path / "roadmap.pptx")
    prs = Presentation(str(path))

    assert path.exists()
    assert len(prs.slides) == 3
    assert any("Roadmap" in shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))


def test_exporter_supports_table_and_template_brand_tokens(tmp_path) -> None:
    outline = OutlineBuilder().build(
        mode=DeckMode.TABLE_TO_DECK,
        title="KPI",
        slide_count=4,
        table_insights={"key_findings": ["Revenue grew"]},
    )
    design = DesignBuilder().build(
        outline=outline,
        brand_tokens={
            "primary_color": "#111111",
            "secondary_color": "#222222",
            "accent_color": "#333333",
            "font_heading": "Aptos Display",
            "font_body": "Aptos",
        },
    )
    ir = SlideIrBuilder().build(
        outline=outline,
        spec_lock=design["spec_lock"],
        table_insights={"key_findings": ["Revenue grew"]},
        chart_specs=[{"type": "line", "title": "Revenue trend"}],
    )

    path = PptxExporter().export(ir, tmp_path / "kpi.pptx")

    assert Presentation(str(path)).slides


def _minimal_content_for(slide_type: SlideType) -> dict:
    """Smallest payload that keeps each renderer happy."""
    if slide_type == SlideType.COVER:
        return {"title": "Smoke", "subtitle": "Test"}
    if slide_type == SlideType.SECTION:
        return {"section_title": "Phase 1"}
    if slide_type == SlideType.CLOSING:
        return {"headline": "Thank you", "body": "Q&A"}
    if slide_type == SlideType.AGENDA:
        return {"items": ["A", "B", "C"]}
    if slide_type == SlideType.CONTENT:
        return {"body": "Body sentence.", "bullets": ["x", "y"]}
    if slide_type == SlideType.COMPARISON:
        return {
            "left": {"title": "L", "bullets": ["a"]},
            "right": {"title": "R", "bullets": ["b"]},
        }
    if slide_type == SlideType.TIMELINE:
        return {
            "milestones": [
                {"label": "Q1", "title": "Plan", "description": "kickoff"},
                {"label": "Q2", "title": "Pilot", "description": "validate"},
            ]
        }
    if slide_type == SlideType.METRIC_CARDS:
        return {"metrics": [{"label": "Adoption", "value": "+18%"}]}
    if slide_type in {SlideType.CHART_BAR, SlideType.CHART_LINE, SlideType.CHART_PIE}:
        return {
            "chart_title": "Trend",
            "categories": ["Q1", "Q2"],
            "series": [{"name": "Revenue", "values": [10.0, 12.0]}],
        }
    if slide_type == SlideType.DATA_TABLE:
        return {"headers": ["Metric", "Value"], "rows": [["Adoption", "18%"]]}
    if slide_type == SlideType.DATA_OVERVIEW:
        return {"body": "Q1 dataset", "bullets": ["clean"], "metrics": []}
    if slide_type == SlideType.INSIGHT_SUMMARY:
        return {"findings": ["Adoption rose"], "risks": ["Sample size"]}
    if slide_type == SlideType.SUMMARY:
        return {"body": "Wrap up", "bullets": ["next step"]}
    return {"body": "x", "bullets": ["x"]}


def test_exporter_renders_every_slide_type_without_raising(tmp_path) -> None:
    """Smoke: each registered SlideType produces a slide using minimal content."""
    slides = []
    for index, slide_type in enumerate(LAYOUT_REGISTRY.keys(), start=1):
        slides.append(
            {
                "id": f"slide_{index:02d}",
                "index": index,
                "title": f"{slide_type.value} title",
                "slide_type": slide_type.value,
                "content": _minimal_content_for(slide_type),
                "theme": {
                    "primary_color": "#0EA5E9",
                    "secondary_color": "#1E293B",
                    "accent_color": "#FACC15",
                    "background_color": "#FFFFFF",
                },
                "assets": {},
            }
        )
    ir = {
        "version": 1,
        "mode": DeckMode.TOPIC_TO_DECK.value,
        "title": "All layouts",
        "slides": slides,
    }
    path = PptxExporter().export(ir, tmp_path / "all_layouts.pptx")
    prs = Presentation(str(path))
    assert len(prs.slides) == len(LAYOUT_REGISTRY)


def test_exporter_consumes_resolved_assets(tmp_path) -> None:
    """When a renderer is given an image_path / icon, no exception is raised."""
    image_path = tmp_path / "fake_asset.png"
    # Smallest valid PNG (1×1 transparent pixel)
    image_path.write_bytes(
        bytes.fromhex(
            "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
            "890000000d49444154789c63000100000005000100"
            "0d0a2db40000000049454e44ae426082"
        )
    )

    slides = [
        {
            "id": "slide_01",
            "index": 1,
            "title": "Cover",
            "slide_type": "cover",
            "content": {"title": "Hello", "subtitle": "World"},
            "theme": {
                "primary_color": "#0EA5E9",
                "secondary_color": "#1E293B",
                "accent_color": "#FACC15",
                "background_color": "#FFFFFF",
            },
            "assets": {"image_path": str(image_path)},
        }
    ]
    ir = {"version": 1, "title": "Cover", "slides": slides}
    path = PptxExporter().export(ir, tmp_path / "cover.pptx")
    assert path.exists()


def test_exporter_rejects_empty_deck(tmp_path) -> None:
    from ppt_exporter import PptxExportError

    with pytest.raises(PptxExportError):
        PptxExporter().export({"slides": []}, tmp_path / "empty.pptx")


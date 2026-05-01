"""Editable PPTX export for ppt-maker.

Each slide IR entry is dispatched to a layout-specific renderer that knows the
schema produced by ``ppt_layouts.LAYOUT_REGISTRY``. The exporter never raises
on missing optional fields — it always produces a renderable slide so the audit
step can flag empty zones rather than failing the whole export.
"""

from __future__ import annotations

import logging
from pathlib import Path
from types import SimpleNamespace
from typing import Any, Callable

Presentation: Any = None
CategoryChartData: Any = None
RGBColor: Any = None
XL_CHART_TYPE: Any = None
MSO_SHAPE: Any = None
PP_ALIGN: Any = SimpleNamespace(LEFT=None, CENTER=None)
Emu: Any = None
Inches: Any = None
Pt: Any = None

logger = logging.getLogger(__name__)

SLIDE_W: Any = None
SLIDE_H: Any = None


class PptxExportError(RuntimeError):
    """Raised when PPTX export fails."""


def _load_pptx() -> None:
    """Import python-pptx only when editable export is actually requested."""
    global CategoryChartData, Emu, Inches, MSO_SHAPE, PP_ALIGN, Presentation, Pt
    global RGBColor, SLIDE_H, SLIDE_W, XL_CHART_TYPE
    if Presentation is not None:
        return
    try:
        from pptx import Presentation as _Presentation  # type: ignore[import-not-found]
        from pptx.chart.data import CategoryChartData as _CategoryChartData  # type: ignore[import-not-found]
        from pptx.dml.color import RGBColor as _RGBColor  # type: ignore[import-not-found]
        from pptx.enum.chart import XL_CHART_TYPE as _XL_CHART_TYPE  # type: ignore[import-not-found]
        from pptx.enum.shapes import MSO_SHAPE as _MSO_SHAPE  # type: ignore[import-not-found]
        from pptx.enum.text import PP_ALIGN as _PP_ALIGN  # type: ignore[import-not-found]
        from pptx.util import Emu as _Emu, Inches as _Inches, Pt as _Pt  # type: ignore[import-not-found]
    except ModuleNotFoundError as exc:
        if exc.name == "pptx":
            raise PptxExportError(
                "Missing optional dependency 'python-pptx'. "
                "Install the ppt-maker advanced_export dependency before exporting PPTX."
            ) from exc
        raise

    Presentation = _Presentation
    CategoryChartData = _CategoryChartData
    RGBColor = _RGBColor
    XL_CHART_TYPE = _XL_CHART_TYPE
    MSO_SHAPE = _MSO_SHAPE
    PP_ALIGN = _PP_ALIGN
    Emu = _Emu
    Inches = _Inches
    Pt = _Pt
    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)


class PptxExporter:
    """Render slide IR into editable PowerPoint shapes."""

    def export(self, slides_ir: dict[str, Any], output_path: str | Path) -> Path:
        _load_pptx()
        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
        theme = self._theme(slides_ir)
        for slide_ir in slides_ir.get("slides", []):
            self._render_slide(prs, slide_ir, theme)
        if len(prs.slides) == 0:
            raise PptxExportError("Cannot export an empty deck")
        prs.save(path)
        return path

    # ── Dispatcher ─────────────────────────────────────────────────────

    def _render_slide(self, prs: Presentation, slide_ir: dict[str, Any], theme: dict[str, str]) -> None:
        slide_type = slide_ir.get("slide_type", "content")
        renderer = self._renderers().get(slide_type, self._render_content)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        try:
            renderer(slide, slide_ir, theme)
        except Exception as exc:  # noqa: BLE001
            logger.warning("ppt-maker: renderer failed for %s: %s", slide_type, exc)
            # Last-resort: render whatever bullets/body we have.
            self._render_content(slide, slide_ir, theme)

    def _renderers(self) -> dict[str, Callable[[Any, dict[str, Any], dict[str, str]], None]]:
        return {
            "cover": self._render_cover,
            "section": self._render_section,
            "agenda": self._render_agenda,
            "content": self._render_content,
            "comparison": self._render_comparison,
            "timeline": self._render_timeline,
            "metric_cards": self._render_metric_cards,
            "chart_bar": self._render_chart,
            "chart_line": self._render_chart,
            "chart_pie": self._render_chart,
            "data_table": self._render_data_table,
            "data_overview": self._render_data_overview,
            "insight_summary": self._render_insight_summary,
            "summary": self._render_summary,
            "closing": self._render_closing,
        }

    # ── Renderers ──────────────────────────────────────────────────────

    def _render_cover(self, slide: Any, ir: dict[str, Any], theme: dict[str, str]) -> None:
        content = ir.get("content") or {}
        # Background accent bar
        self._add_bar(slide, theme["primary_color"], top=3.0, height=0.15)
        title = content.get("title") or ir.get("title") or "Untitled"
        subtitle = content.get("subtitle") or ""
        self._textbox(slide, title, left=0.8, top=2.2, width=11.7, height=1.4,
                      size=44, bold=True, color=theme["primary_color"])
        if subtitle:
            self._textbox(slide, subtitle, left=0.8, top=3.6, width=11.7, height=1.2,
                          size=22, color=theme["secondary_color"])
        meta = " · ".join(filter(None, [content.get("presenter") or "", content.get("date") or ""]))
        if meta:
            self._textbox(slide, meta, left=0.8, top=6.2, width=11.7, height=0.6,
                          size=14, color=theme["secondary_color"])
        self._maybe_image(slide, ir, left=Inches(8.5), top=Inches(0.5), width=Inches(4.3), height=Inches(1.6))

    def _render_section(self, slide: Any, ir: dict[str, Any], theme: dict[str, str]) -> None:
        content = ir.get("content") or {}
        self._add_full_band(slide, theme["secondary_color"])
        title = content.get("section_title") or ir.get("title", "")
        subtitle = content.get("section_subtitle") or ""
        self._textbox(slide, title, left=0.8, top=2.6, width=11.7, height=1.6,
                      size=44, bold=True, color="#FFFFFF")
        if subtitle:
            self._textbox(slide, subtitle, left=0.8, top=4.4, width=11.7, height=1.0,
                          size=20, color="#FFFFFF")

    def _render_agenda(self, slide: Any, ir: dict[str, Any], theme: dict[str, str]) -> None:
        self._title(slide, ir.get("title", "Agenda"), theme)
        items = (ir.get("content") or {}).get("items") or []
        items = [it for it in items if it]
        if not items:
            items = ["背景与目标", "现状与挑战", "核心方案", "总结与下一步"]
        per_col = max(1, (len(items) + 1) // 2)
        for index, item in enumerate(items[:10]):
            col = 0 if index < per_col else 1
            row = index if col == 0 else index - per_col
            left = Inches(0.9 + col * 6.3)
            top = Inches(1.6 + row * 0.85)
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.55), Inches(0.55))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(theme["primary_color"])
            shape.line.color.rgb = self._rgb(theme["primary_color"])
            shape.text_frame.text = str(index + 1)
            shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            shape.text_frame.paragraphs[0].font.bold = True
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            self._textbox(slide, item, left=left + Inches(0.7), top=top, width=Inches(5.4),
                          height=Inches(0.55), size=18, color=theme["secondary_color"])

    def _render_content(self, slide: Any, ir: dict[str, Any], theme: dict[str, str]) -> None:
        self._title(slide, ir.get("title", ""), theme)
        content = ir.get("content") or {}
        body = content.get("body") or ir.get("notes") or ""
        bullets = [b for b in (content.get("bullets") or []) if b]
        has_image = bool((ir.get("assets") or {}).get("image_path"))
        body_width = 7.6 if has_image else 11.8
        if body:
            self._textbox(slide, body, left=0.8, top=1.4, width=body_width, height=1.0,
                          size=18, color=theme["secondary_color"])
        if bullets:
            self._bullets(slide, bullets, left=0.8, top=2.5, width=body_width, theme=theme)
        if has_image:
            self._maybe_image(slide, ir, left=Inches(8.6), top=Inches(1.4), width=Inches(4.0), height=Inches(4.5))
        else:
            self._maybe_icon(slide, ir, left=Inches(11.6), top=Inches(1.4), size=Inches(1.0), theme=theme)

    def _render_comparison(self, slide: Any, ir: dict[str, Any], theme: dict[str, str]) -> None:
        self._title(slide, ir.get("title", ""), theme)
        content = ir.get("content") or {}
        left = content.get("left") or {}
        right = content.get("right") or {}
        for index, side in enumerate((left, right)):
            base_left = 0.8 + index * 6.2
            self._add_panel(slide, theme["primary_color"] if index == 0 else theme["accent_color"],
                            left=base_left, top=1.4, width=5.85, height=4.6)
            self._textbox(slide, side.get("title") or ("现状" if index == 0 else "新方案"),
                          left=base_left + 0.2, top=1.55, width=5.45, height=0.6,
                          size=22, bold=True, color="#FFFFFF")
            bullets = side.get("bullets") or []
            self._bullets(slide, bullets, left=base_left + 0.2, top=2.3, width=5.45,
                          theme=theme, color="#FFFFFF")
        summary = content.get("summary")
        if summary:
            self._textbox(slide, summary, left=0.8, top=6.3, width=11.7, height=0.7,
                          size=16, color=theme["secondary_color"])

    def _render_timeline(self, slide: Any, ir: dict[str, Any], theme: dict[str, str]) -> None:
        self._title(slide, ir.get("title", ""), theme)
        milestones = (ir.get("content") or {}).get("milestones") or []
        if not milestones:
            return
        track_top = Inches(3.5)
        slide.shapes.add_connector(1, Inches(0.8), track_top + Emu(50000), Inches(12.5),
                                    track_top + Emu(50000))
        gap = max(1.2, (12.0 / max(1, len(milestones[:6]))))
        for index, milestone in enumerate(milestones[:6]):
            cx = Inches(0.8 + index * gap)
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Emu(80000), track_top - Emu(80000),
                                          Inches(0.35), Inches(0.35))
            dot.fill.solid()
            dot.fill.fore_color.rgb = self._rgb(theme["primary_color"])
            dot.line.color.rgb = self._rgb(theme["primary_color"])
            self._textbox(slide, milestone.get("label", f"M{index + 1}"),
                          left=cx - Inches(0.6), top=Inches(2.7), width=Inches(2.0), height=Inches(0.5),
                          size=14, bold=True, color=theme["primary_color"])
            self._textbox(slide, milestone.get("title", ""),
                          left=cx - Inches(0.8), top=Inches(4.0), width=Inches(2.4), height=Inches(0.7),
                          size=16, bold=True, color=theme["secondary_color"])
            description = milestone.get("description")
            if description:
                self._textbox(slide, description,
                              left=cx - Inches(0.8), top=Inches(4.7), width=Inches(2.4), height=Inches(1.5),
                              size=12, color=theme["secondary_color"])

    def _render_metric_cards(self, slide: Any, ir: dict[str, Any], theme: dict[str, str]) -> None:
        self._title(slide, ir.get("title", ""), theme)
        metrics = (ir.get("content") or {}).get("metrics") or []
        if not metrics:
            return
        count = min(len(metrics), 4)
        card_w = (12.0 - (count - 1) * 0.3) / count
        for index, metric in enumerate(metrics[:count]):
            left = 0.8 + index * (card_w + 0.3)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(1.6), Inches(card_w), Inches(2.4),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(theme["primary_color"])
            shape.line.color.rgb = self._rgb(theme["primary_color"])
            frame = shape.text_frame
            frame.clear()
            label = frame.paragraphs[0]
            label.text = str(metric.get("label", ""))
            label.font.size = Pt(14)
            label.font.color.rgb = RGBColor(255, 255, 255)
            label.alignment = PP_ALIGN.CENTER
            value = frame.add_paragraph()
            value.text = str(metric.get("value", "—"))
            value.font.size = Pt(36)
            value.font.bold = True
            value.font.color.rgb = RGBColor(255, 255, 255)
            value.alignment = PP_ALIGN.CENTER
            delta = metric.get("delta")
            if delta:
                d = frame.add_paragraph()
                d.text = str(delta)
                d.font.size = Pt(12)
                d.font.color.rgb = RGBColor(255, 255, 255)
                d.alignment = PP_ALIGN.CENTER
        bullets = (ir.get("content") or {}).get("bullets") or []
        if bullets:
            self._bullets(slide, bullets, left=0.8, top=4.4, width=11.8, theme=theme)

    def _render_chart(self, slide: Any, ir: dict[str, Any], theme: dict[str, str]) -> None:
        self._title(slide, ir.get("title", ""), theme)
        content = ir.get("content") or {}
        categories = content.get("categories") or []
        series = content.get("series") or []
        chart_type_str = (content.get("chart_type") or ir.get("slide_type", "chart_bar").replace("chart_", "")).lower()
        chart_data = CategoryChartData()
        if categories and series:
            chart_data.categories = list(categories)
            for s in series[:5]:
                values = list(s.get("values") or [])
                if not values:
                    continue
                chart_data.add_series(str(s.get("name") or "系列"), tuple(values))
        else:
            chart_data.categories = ["A", "B", "C", "D"]
            chart_data.add_series("示意", (40, 65, 50, 80))
        try:
            chart_type = {
                "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "horizontal_bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "line": XL_CHART_TYPE.LINE,
                "pie": XL_CHART_TYPE.PIE,
            }.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)
            slide.shapes.add_chart(
                chart_type,
                Inches(0.8), Inches(1.4), Inches(8.5), Inches(4.8),
                chart_data,
            )
        except Exception as exc:  # noqa: BLE001
            logger.info("ppt-maker chart fallback: %s", exc)
            self._textbox(slide, content.get("chart_title", ""), left=0.8, top=2.0,
                          width=11.7, height=4.0, size=20, color=theme["primary_color"])
        bullets = content.get("bullets") or []
        if bullets:
            self._bullets(slide, bullets, left=9.6, top=1.6, width=3.2, theme=theme)

    def _render_data_table(self, slide: Any, ir: dict[str, Any], theme: dict[str, str]) -> None:
        self._title(slide, ir.get("title", ""), theme)
        content = ir.get("content") or {}
        headers = list(content.get("headers") or ["维度", "数值"])[:8]
        rows = [list(r) for r in (content.get("rows") or [])][:10]
        if not rows:
            rows = [[""] * len(headers)]
        rendered = slide.shapes.add_table(
            len(rows) + 1,
            len(headers),
            Inches(0.7), Inches(1.4), Inches(11.9), Inches(4.6),
        ).table
        for col_index, header in enumerate(headers):
            cell = rendered.cell(0, col_index)
            cell.text = str(header)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._rgb(theme["primary_color"])
        for r_index, row in enumerate(rows, start=1):
            for c_index, value in enumerate(row[: len(headers)]):
                rendered.cell(r_index, c_index).text = str(value)
        bullets = content.get("bullets") or []
        if bullets:
            self._bullets(slide, bullets, left=0.7, top=6.2, width=11.9, theme=theme)

    def _render_data_overview(self, slide: Any, ir: dict[str, Any], theme: dict[str, str]) -> None:
        self._title(slide, ir.get("title", ""), theme)
        content = ir.get("content") or {}
        body = content.get("body") or ""
        bullets = content.get("bullets") or []
        if body:
            self._textbox(slide, body, left=0.8, top=1.4, width=11.7, height=1.0,
                          size=18, color=theme["secondary_color"])
        if bullets:
            self._bullets(slide, bullets, left=0.8, top=2.5, width=7.5, theme=theme)
        metrics = content.get("metrics") or []
        for index, metric in enumerate(metrics[:3]):
            left = 9.0 + index * 1.3
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(2.5), Inches(1.2), Inches(2.0),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(theme["accent_color"])
            shape.line.color.rgb = self._rgb(theme["accent_color"])
            shape.text_frame.text = str(metric.get("value", ""))
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _render_insight_summary(self, slide: Any, ir: dict[str, Any], theme: dict[str, str]) -> None:
        self._title(slide, ir.get("title", ""), theme)
        content = ir.get("content") or {}
        findings = content.get("findings") or []
        risks = content.get("risks") or []
        self._textbox(slide, "核心发现", left=0.8, top=1.4, width=5.7, height=0.5,
                      size=18, bold=True, color=theme["primary_color"])
        self._bullets(slide, findings, left=0.8, top=2.0, width=5.7, theme=theme)
        self._textbox(slide, "风险提示", left=7.0, top=1.4, width=5.5, height=0.5,
                      size=18, bold=True, color=theme["accent_color"])
        self._bullets(slide, risks, left=7.0, top=2.0, width=5.5, theme=theme)

    def _render_summary(self, slide: Any, ir: dict[str, Any], theme: dict[str, str]) -> None:
        self._title(slide, ir.get("title", "总结"), theme)
        content = ir.get("content") or {}
        body = content.get("body") or ""
        bullets = content.get("bullets") or []
        if body:
            self._textbox(slide, body, left=0.8, top=1.4, width=11.7, height=1.0,
                          size=20, color=theme["secondary_color"])
        if bullets:
            self._bullets(slide, bullets, left=0.8, top=2.6, width=11.7, theme=theme)

    def _render_closing(self, slide: Any, ir: dict[str, Any], theme: dict[str, str]) -> None:
        content = ir.get("content") or {}
        self._add_full_band(slide, theme["primary_color"])
        headline = content.get("headline") or ir.get("title") or "Thank You"
        body = content.get("body") or ""
        contact = content.get("contact") or ""
        self._textbox(slide, headline, left=0.8, top=2.6, width=11.7, height=1.6,
                      size=54, bold=True, color="#FFFFFF")
        if body:
            self._textbox(slide, body, left=0.8, top=4.4, width=11.7, height=1.2,
                          size=20, color="#FFFFFF")
        if contact:
            self._textbox(slide, contact, left=0.8, top=6.0, width=11.7, height=0.7,
                          size=16, color="#FFFFFF")

    # ── Primitives ─────────────────────────────────────────────────────

    def _title(self, slide: Any, text: str, theme: dict[str, str]) -> None:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(12.0), Inches(0.8))
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = text or ""
        paragraph.font.size = Pt(28)
        paragraph.font.bold = True
        paragraph.font.name = theme["font_heading"]
        paragraph.font.color.rgb = self._rgb(theme["primary_color"])
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(0.6), Inches(0.06),
        ).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = self._rgb(theme["accent_color"])
        slide.shapes[-1].line.fill.background()

    def _textbox(
        self,
        slide: Any,
        text: str,
        *,
        left: float | Any,
        top: float | Any,
        width: float | Any,
        height: float | Any,
        size: int = 16,
        bold: bool = False,
        color: str | None = None,
        align: Any = None,
    ) -> None:
        align = PP_ALIGN.LEFT if align is None else align
        l = left if not isinstance(left, (int, float)) else Inches(left)
        t = top if not isinstance(top, (int, float)) else Inches(top)
        w = width if not isinstance(width, (int, float)) else Inches(width)
        h = height if not isinstance(height, (int, float)) else Inches(height)
        box = slide.shapes.add_textbox(l, t, w, h)
        frame = box.text_frame
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.text = text or ""
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.alignment = align
        if color:
            paragraph.font.color.rgb = self._rgb(color)

    def _bullets(
        self,
        slide: Any,
        bullets: list[str],
        *,
        left: float = 0.8,
        top: float = 1.45,
        width: float = 11.7,
        theme: dict[str, str],
        title: str | None = None,
        color: str | None = None,
    ) -> None:
        if title:
            self._textbox(slide, title, left=left, top=top, width=width, height=0.4,
                          size=18, bold=True, color=color or theme["primary_color"])
            top += 0.5
        items = [str(item) for item in bullets if item]
        if not items:
            return
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.0))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        for index, bullet in enumerate(items):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = f"• {bullet}"
            paragraph.font.size = Pt(16)
            paragraph.font.name = theme["font_body"]
            if color:
                paragraph.font.color.rgb = self._rgb(color)
            paragraph.space_after = Pt(6)

    def _add_bar(self, slide: Any, color: str, *, top: float, height: float) -> None:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(top), SLIDE_W, Inches(height),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb(color)
        bar.line.fill.background()

    def _add_full_band(self, slide: Any, color: str) -> None:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb(color)
        bar.line.fill.background()

    def _add_panel(self, slide: Any, color: str, *, left: float, top: float,
                   width: float, height: float) -> None:
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = self._rgb(color)
        panel.line.fill.background()

    def _maybe_image(
        self,
        slide: Any,
        ir: dict[str, Any],
        *,
        left: Any,
        top: Any,
        width: Any,
        height: Any,
    ) -> None:
        path = ((ir.get("assets") or {}).get("image_path") or "").strip()
        if not path:
            return
        try:
            slide.shapes.add_picture(path, left, top, width=width, height=height)
        except Exception as exc:  # noqa: BLE001
            logger.info("ppt-maker image insert failed (%s): %s", path, exc)

    def _maybe_icon(
        self,
        slide: Any,
        ir: dict[str, Any],
        *,
        left: Any,
        top: Any,
        size: Any,
        theme: dict[str, str],
    ) -> None:
        icon = (ir.get("assets") or {}).get("icon")
        if not icon:
            return
        try:
            shape_enum = icon.get("shape")
            if isinstance(shape_enum, str):
                shape_enum = getattr(MSO_SHAPE, shape_enum, MSO_SHAPE.OVAL)
            shape = slide.shapes.add_shape(shape_enum or MSO_SHAPE.OVAL, left, top, size, size)
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(theme["accent_color"])
            shape.line.color.rgb = self._rgb(theme["accent_color"])
            emoji = icon.get("emoji")
            if emoji:
                shape.text_frame.text = str(emoji)
                shape.text_frame.paragraphs[0].font.size = Pt(18)
                shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        except Exception as exc:  # noqa: BLE001
            logger.info("ppt-maker icon insert failed: %s", exc)

    def _theme(self, slides_ir: dict[str, Any]) -> dict[str, str]:
        first = next(iter(slides_ir.get("slides", [])), {})
        theme = first.get("theme") or {}
        return {
            "primary_color": theme.get("primary_color") or "#3457D5",
            "secondary_color": theme.get("secondary_color") or "#172033",
            "accent_color": theme.get("accent_color") or "#FFB000",
            "font_heading": theme.get("font_heading") or "Microsoft YaHei",
            "font_body": theme.get("font_body") or "Microsoft YaHei",
        }

    @staticmethod
    def _rgb(color: str) -> RGBColor:
        cleaned = (color or "").lstrip("#")
        if len(cleaned) != 6:
            cleaned = "3457D5"
        try:
            return RGBColor(int(cleaned[0:2], 16), int(cleaned[2:4], 16), int(cleaned[4:6], 16))
        except ValueError:
            return RGBColor(0x34, 0x57, 0xD5)
